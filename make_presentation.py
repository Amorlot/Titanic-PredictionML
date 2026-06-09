from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x23, 0x4E)
BLUE       = RGBColor(0x23, 0x5A, 0xA3)
TEAL       = RGBColor(0x17, 0x8A, 0x8A)
ORANGE     = RGBColor(0xE8, 0x7B, 0x1E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
DARK_GRAY  = RGBColor(0x2C, 0x2C, 0x2C)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xC0, 0x39, 0x2B)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)
DARK_BG    = RGBColor(0x1E, 0x1E, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size=13, bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.1), fill=NAVY)
    txt(slide, title, Inches(0.5), Inches(0.1), W - Inches(1), Inches(0.65),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.72), W - Inches(1), Inches(0.36),
            size=14, color=RGBColor(0xAA, 0xC4, 0xE8), italic=True)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.05), fill=ORANGE)


def pnum(slide, n):
    txt(slide, str(n), W - Inches(0.55), H - Inches(0.38),
        Inches(0.4), Inches(0.3), size=11,
        color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# S1 — TITLE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(2.85), W, Inches(0.07), fill=ORANGE)
add_rect(sl, 0, Inches(5.3), W, Inches(0.07), fill=TEAL)

txt(sl, "TITANIC SURVIVAL PREDICTION",
    Inches(0.8), Inches(1.1), W - Inches(1.6), Inches(1.05),
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Machine Learning Pipeline — Architettura Avanzata con Libreria Generalizzata",
    Inches(0.8), Inches(2.15), W - Inches(1.6), Inches(0.6),
    size=19, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER, italic=True)
txt(sl, "5 Algoritmi  |  Docker  |  REST API  |  Feature Engineering  |  Libreria Riutilizzabile",
    Inches(0.8), Inches(3.1), W - Inches(1.6), Inches(0.45),
    size=13, color=RGBColor(0x78, 0xA0, 0xCC), align=PP_ALIGN.CENTER)

stats = [("5", "Modelli"), ("891", "Train rows"),
         ("418", "Test rows"), ("84.4%", "Best CV acc.")]
for i, (val, lab) in enumerate(stats):
    bx = Inches(1.5) + i * Inches(2.6)
    add_rect(sl, bx, Inches(5.55), Inches(2.2), Inches(1.45),
             fill=RGBColor(0x2A, 0x3A, 0x6E), line=BLUE, line_w=Pt(1.5))
    txt(sl, val, bx, Inches(5.65), Inches(2.2), Inches(0.65),
        size=34, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, lab, bx, Inches(6.32), Inches(2.2), Inches(0.42),
        size=13, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

txt(sl, "Corso Data Analytics  |  2025/26  |  Andrea Morlotti",
    Inches(0.5), H - Inches(0.5), W - Inches(1.0), Inches(0.35),
    size=11, color=RGBColor(0x55, 0x66, 0x88), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S2 — INDICE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Indice", "Struttura della presentazione")
pnum(sl, 2)

sections = [
    ("01", "Il Problema",             "Obiettivo Kaggle e analisi del dataset"),
    ("02", "Architettura",            "Docker, REST API, config.yaml"),
    ("03", "Libreria Generalizzata",  "Abstract classes, composabilita, riuso"),
    ("04", "Data Pipeline",           "Sequenza completa dei transformer"),
    ("05", "Feature Engineering",     "7 nuove variabili derivate"),
    ("06", "Imputazione Age",         "Strategia per titolo del nome"),
    ("07", "Modelli & GridSearchCV",  "5 algoritmi con ottimizzazione"),
    ("08", "Tecniche Avanzate",       "PCA, Multicollinearity, Threshold Tuning"),
    ("09", "Parallelismo & Docker",   "Fix deadlock, LOKY, n_jobs"),
    ("10", "Valutazione Locale",      "StratifiedKFold vs holdout"),
    ("11", "Risultati",               "Performance CV per modello"),
    ("12", "Confronto Funghi-Titanic","Evoluzione architettura"),
    ("13", "Prossimi Passi",          "Ensemble, CV-Bagging, roadmap"),
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 4
    col = i % 4
    bx = Inches(0.28) + col * Inches(3.2)
    by = Inches(1.3)  + row * Inches(1.9)
    add_rect(sl, bx, by, Inches(3.0), Inches(1.72),
             fill=LIGHT_GRAY, line=BLUE, line_w=Pt(1))
    add_rect(sl, bx, by, Inches(0.5), Inches(1.72), fill=BLUE)
    txt(sl, num, bx, by + Inches(0.62), Inches(0.5), Inches(0.5),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, bx + Inches(0.58), by + Inches(0.14),
        Inches(2.3), Inches(0.5), size=12, bold=True, color=NAVY)
    txt(sl, desc, bx + Inches(0.58), by + Inches(0.66),
        Inches(2.3), Inches(0.85), size=10, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# S3 — IL PROBLEMA
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Il Problema — Titanic Kaggle", "Binary classification: Survived = 0 / 1")
pnum(sl, 3)

txt(sl, "Obiettivo", Inches(0.4), Inches(1.28), Inches(5.8), Inches(0.38),
    size=15, bold=True, color=NAVY)
txt(sl,
    "Predire la sopravvivenza di 418 passeggeri del test set in base a caratteristiche "
    "demografiche e sociali. Kaggle valuta con Accuracy (% predizioni corrette).",
    Inches(0.4), Inches(1.68), Inches(5.8), Inches(0.85),
    size=12, color=DARK_GRAY)

add_rect(sl, Inches(0.4), Inches(2.6), Inches(5.8), Inches(2.5), fill=LIGHT_GRAY)
add_rect(sl, Inches(0.4), Inches(2.6), Inches(5.8), Inches(0.38), fill=BLUE)
txt(sl, "Dataset Kaggle", Inches(0.54), Inches(2.62),
    Inches(5.6), Inches(0.34), size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "891 passeggeri nel train set (labels noti)",
    "418 passeggeri nel test set (labels segreti)",
    "Metrica Kaggle: Accuracy",
    "Baseline naive (tutti morti): 61.6%",
    "Obiettivo realistico: >80%",
]):
    txt(sl, f"  {item}", Inches(0.54), Inches(3.06) + i * Inches(0.37),
        Inches(5.6), Inches(0.34), size=12, color=DARK_GRAY)

add_rect(sl, Inches(0.4), Inches(5.2), Inches(5.8), Inches(2.05),
         fill=RGBColor(0xE8, 0xF4, 0xF8))
add_rect(sl, Inches(0.4), Inches(5.2), Inches(5.8), Inches(0.38),
         fill=RGBColor(0x17, 0x8A, 0x8A))
txt(sl, "Feature originali (11 colonne)", Inches(0.54), Inches(5.22),
    Inches(5.6), Inches(0.34), size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "PassengerId, Pclass (1/2/3), Name",
    "Sex, Age (177 mancanti), SibSp, Parch",
    "Ticket, Fare, Cabin (687 mancanti), Embarked",
]):
    txt(sl, f"  {item}", Inches(0.54), Inches(5.66) + i * Inches(0.37),
        Inches(5.6), Inches(0.34), size=12, color=DARK_GRAY)

# Barre sopravvivenza
add_rect(sl, Inches(6.6), Inches(1.28), Inches(6.3), Inches(5.9),
         fill=RGBColor(0xF8, 0xF9, 0xFF), line=NAVY, line_w=Pt(1.5))
add_rect(sl, Inches(6.6), Inches(1.28), Inches(6.3), Inches(0.42), fill=NAVY)
txt(sl, "Tassi di sopravvivenza (train set)",
    Inches(6.72), Inches(1.3), Inches(6.1), Inches(0.38),
    size=13, bold=True, color=WHITE)

bars = [
    ("Donne",              74.2, GREEN),
    ("1a classe",          62.9, GREEN),
    ("Bambini < 15 anni",  57.9, TEAL),
    ("Totale passeggeri",  38.4, ORANGE),
    ("3a classe",          24.2, RED),
    ("Uomini",             18.9, RED),
]
for i, (label, pct, col) in enumerate(bars):
    y = Inches(1.9) + i * Inches(0.72)
    txt(sl, label, Inches(6.75), y, Inches(2.0), Inches(0.3), size=11, color=DARK_GRAY)
    add_rect(sl, Inches(8.85), y + Inches(0.05), Inches(3.75), Inches(0.2),
             fill=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(sl, Inches(8.85), y + Inches(0.05),
             pct / 100 * Inches(3.75), Inches(0.2), fill=col)
    txt(sl, f"{pct}%", Inches(12.68), y, Inches(0.72), Inches(0.3),
        size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# S4 — ARCHITETTURA
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Architettura del Progetto", "Docker  |  REST API  |  config.yaml  |  Libreria embedded")
pnum(sl, 4)

add_rect(sl, Inches(0.3), Inches(1.28), Inches(4.0), Inches(6.0), fill=DARK_BG)
txt(sl, "titanic/", Inches(0.45), Inches(1.4), Inches(3.7), Inches(0.32),
    size=12, bold=True, color=ORANGE)
tree = [
    ("|- lib/",                           TEAL,                        12),
    ("|   |- models/",                    RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |   |- base.py",                WHITE,                       10),
    ("|   |   |- logreg.py",              WHITE,                       10),
    ("|   |   |- random_forest.py",       WHITE,                       10),
    ("|   |   |- xgboost.py",             WHITE,                       10),
    ("|   |   |- svm.py",                 WHITE,                       10),
    ("|   |   +- decision_tree.py",       WHITE,                       10),
    ("|   |- api/models.py",              RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |- feature_engineering.py",     RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |- cleaner.py",                 RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |- encoder.py",                 RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |- outliers.py",                RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   |- multicollinearity.py",       RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|   +- pca_reducer.py",             RGBColor(0xAA, 0xCC, 0xFF), 11),
    ("|- main.py",                        ORANGE,                      12),
    ("|- config.yaml",                    TEAL,                        12),
    ("|- Dockerfile",                     RGBColor(0xFF, 0xCC, 0x44), 12),
    ("|- evaluate_local.py",              GREEN,                       12),
    ("+- tests/",                         RGBColor(0xBB, 0xBB, 0xFF), 12),
]
for i, (line, col, sz) in enumerate(tree):
    txt(sl, line, Inches(0.48), Inches(1.76) + i * Inches(0.25),
        Inches(3.7), Inches(0.26), size=sz, color=col)

add_rect(sl, Inches(4.65), Inches(1.28), Inches(4.0), Inches(2.6),
         fill=RGBColor(0xE8, 0xF4, 0xFF), line=BLUE, line_w=Pt(1.5))
add_rect(sl, Inches(4.65), Inches(1.28), Inches(4.0), Inches(0.4), fill=BLUE)
txt(sl, "Docker Container", Inches(4.77), Inches(1.3), Inches(3.8), Inches(0.37),
    size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "FROM python:3.11-slim",
    "ENV PYTHONPATH=/app:/app/lib",
    "ENV PYTHONUNBUFFERED=1",
    "ENV LOKY_MAX_CPU_COUNT=8",
    "RUN apt install gcc  (per XGBoost)",
]):
    txt(sl, f"  {item}", Inches(4.77), Inches(1.78) + i * Inches(0.34),
        Inches(3.75), Inches(0.3), size=11, color=DARK_GRAY)

add_rect(sl, Inches(4.65), Inches(4.0), Inches(4.0), Inches(3.18),
         fill=RGBColor(0xF0, 0xFF, 0xF0), line=GREEN, line_w=Pt(1.5))
add_rect(sl, Inches(4.65), Inches(4.0), Inches(4.0), Inches(0.4), fill=GREEN)
txt(sl, "REST API -- FastAPI", Inches(4.77), Inches(4.02),
    Inches(3.8), Inches(0.37), size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "GET  /           -> lista modelli attivi",
    "POST /train      -> training + best_score",
    "POST /evaluate   -> accuracy, f1, report",
    "POST /predict    -> genera submission CSV",
    "Risposta JSON strutturata",
    "Avvio: bash start.sh nel container",
]):
    txt(sl, f"  {item}", Inches(4.77), Inches(4.5) + i * Inches(0.38),
        Inches(3.75), Inches(0.34), size=11, color=DARK_GRAY)

add_rect(sl, Inches(9.0), Inches(1.28), Inches(4.0), Inches(5.9),
         fill=RGBColor(0xFF, 0xFB, 0xF0), line=ORANGE, line_w=Pt(1.5))
add_rect(sl, Inches(9.0), Inches(1.28), Inches(4.0), Inches(0.4), fill=ORANGE)
txt(sl, "config.yaml -- punto di controllo",
    Inches(9.12), Inches(1.3), Inches(3.8), Inches(0.37),
    size=12, bold=True, color=WHITE)
cfg_lines = [
    ("loader:",                          NAVY,      12, True),
    ("  train_path, target_col",         DARK_GRAY, 11, False),
    ("  drop_cols: [Ticket, Cabin]",     DARK_GRAY, 11, False),
    ("cleaner:",                         NAVY,      12, True),
    ("  num_strategy: knn",              DARK_GRAY, 11, False),
    ("  cat_strategy: knn",              DARK_GRAY, 11, False),
    ("outliers:",                        NAVY,      12, True),
    ("  detection: [iqr, zscore]",       DARK_GRAY, 11, False),
    ("  apply_to: [logreg]",             DARK_GRAY, 11, False),
    ("multicollinearity:",               NAVY,      12, True),
    ("  threshold: 0.85",                DARK_GRAY, 11, False),
    ("  apply_to: [logreg]",             DARK_GRAY, 11, False),
    ("pca:",                             NAVY,      12, True),
    ("  n_components: 0.95",             DARK_GRAY, 11, False),
    ("  apply_to: [logreg]",             DARK_GRAY, 11, False),
    ("models:",                          NAVY,      12, True),
    ("  active: [logreg, xgboost,",      DARK_GRAY, 11, False),
    ("    rf, dt, svm]",                 DARK_GRAY, 11, False),
    ("model_cache:",                     NAVY,      12, True),
    ("  enabled: false",                 DARK_GRAY, 11, False),
]
for i, (line, col, sz, bd) in enumerate(cfg_lines):
    txt(sl, line, Inches(9.15), Inches(1.78) + i * Inches(0.27),
        Inches(3.75), Inches(0.26), size=sz, bold=bd, color=col)


# ═══════════════════════════════════════════════════════════════════
# S5 — LIBRERIA GENERALIZZATA
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Libreria Generalizzata -- lib/",
       "Componenti riutilizzabili su qualsiasi dataset tabellare, non solo Titanic")
pnum(sl, 5)

add_rect(sl, Inches(0.3), Inches(1.28), Inches(12.7), Inches(0.9),
         fill=RGBColor(0xE8, 0xF0, 0xFF), line=BLUE, line_w=Pt(1))
txt(sl,
    "La cartella lib/ non e' specifica per il Titanic -- e' una libreria ML generica. "
    "Ogni componente prende un DataFrame in input, lo trasforma, e restituisce un DataFrame "
    "indipendentemente dal dominio. Per usarla su un nuovo dataset basta cambiare config.yaml.",
    Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.76),
    size=12, color=DARK_GRAY)

# Abstract base model con codice
add_rect(sl, Inches(0.3), Inches(2.35), Inches(6.0), Inches(4.88), fill=DARK_BG)
txt(sl, "models/base.py -- AbstractModel", Inches(0.45), Inches(2.42),
    Inches(5.8), Inches(0.3), size=12, bold=True, color=ORANGE)
code_lines = [
    ("class AbstractModel(ABC):",                TEAL),
    ("    def __init__(self):",                  WHITE),
    ("        self.model        = None",         RGBColor(0xAA, 0xCC, 0xFF)),
    ("        self.threshold    = 0.5",          RGBColor(0xAA, 0xCC, 0xFF)),
    ("        self.best_score   = None",         RGBColor(0xAA, 0xCC, 0xFF)),
    ("        self.best_params  = None",         RGBColor(0xAA, 0xCC, 0xFF)),
    ("",                                          WHITE),
    ("    @abstractmethod",                      ORANGE),
    ("    def train(self, X, y, cv, scoring):",  WHITE),
    ("        ...",                              RGBColor(0x77, 0x77, 0x77)),
    ("",                                          WHITE),
    ("    def predict(self, X):",                TEAL),
    ("        if self.threshold != 0.5:",        WHITE),
    ("            p = model.predict_proba(X)[:,1]", WHITE),
    ("            return (p >= self.threshold)", WHITE),
    ("        return self.model.predict(X)",     WHITE),
    ("",                                          WHITE),
    ("    def tune_threshold(self, X, y):",      TEAL),
    ("        # OOF probs con cross_val_predict", RGBColor(0x77, 0x99, 0x55)),
    ("        # scannerizza soglie 0.30-0.70",   RGBColor(0x77, 0x99, 0x55)),
    ("        # ottimizza F1 / accuracy / recall", RGBColor(0x77, 0x99, 0x55)),
]
for i, (line, col) in enumerate(code_lines):
    txt(sl, line, Inches(0.45), Inches(2.78) + i * Inches(0.24),
        Inches(5.75), Inches(0.25), size=9.5, color=col)

# Componenti della libreria
components = [
    ("GenericCleaner",
     BLUE,
     "Imputa valori mancanti con KNN (num e cat).\n"
     "OrdinalEncoder per cat prima dell'imputer.\n"
     "Salva le colonne al fit(): trasforma sempre\n"
     "con la stessa struttura del training set.\n"
     "Config: num_strategy, cat_strategy, knn_k."),
    ("GenericOutlierHandler",
     TEAL,
     "Rileva outlier con IQR e/o Z-score.\n"
     "Trattamento: winsorization (clip ai pct).\n"
     "Applicabile solo a colonne specifiche.\n"
     "Config: quale modello riceve dati win.\n"
     "Fit su train, transform su test."),
    ("MulticollinearityHandler",
     PURPLE,
     "Calcola correlazione Pearson coppie num.\n"
     "Rimuove una feature per ogni coppia |r|>th.\n"
     "Encoder separato (encoder_mc) per dataset\n"
     "ridotto (colonne diverse = scaler diverso).\n"
     "Config: threshold, apply_to."),
    ("GenericEncoder",
     ORANGE,
     "StandardScaler per numeriche.\n"
     "OneHotEncoder per categoriche.\n"
     "encode_target() per la variabile Y.\n"
     "decode_target() per risalire alle label.\n"
     "Istanza separata per variante mc."),
    ("PCAReducer",
     GREEN,
     "PCA con varianza spiegata configurabile.\n"
     "Istanza separata per ogni modello.\n"
     "Evita mismatch se modelli hanno\n"
     "colonne diverse (es. logreg no Fare).\n"
     "Config: n_components, apply_to."),
    ("FeatureEngineer",
     RED,
     "Crea feature derivate dal dominio.\n"
     "Imputa Age per TitleGroup al fit().\n"
     "Configurabile: drop_raw_fare_for.\n"
     "E' l'unico componente domain-specific;\n"
     "gli altri sono completamente generici."),
]

for i, (name, col, desc) in enumerate(components):
    bx = Inches(6.6) + (i % 3) * Inches(2.25)
    by = Inches(2.35) + (i // 3) * Inches(2.52)
    add_rect(sl, bx, by, Inches(2.1), Inches(2.38),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1.5))
    add_rect(sl, bx, by, Inches(2.1), Inches(0.37), fill=col)
    txt(sl, name, bx + Inches(0.08), by + Inches(0.04),
        Inches(1.95), Inches(0.3), size=11, bold=True, color=WHITE)
    txt(sl, desc, bx + Inches(0.08), by + Inches(0.44),
        Inches(1.96), Inches(1.87), size=9.5, color=DARK_GRAY)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32),
         fill=RGBColor(0xE8, 0xF0, 0xFF), line=BLUE, line_w=Pt(0.75))
txt(sl,
    "Contratto uniforme: configure() -> fit(df_train) -> transform(df)  "
    "-- composabili in pipeline senza coupling tra componenti.",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S6 — DATA PIPELINE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Data Pipeline", "Dal CSV grezzo ai vettori numerici pronti per i modelli")
pnum(sl, 6)

steps = [
    ("LOADER",        "Legge CSV\nDrop cols\nSepara X, y",                NAVY),
    ("FEATURE ENG.",  "7 nuove feature\nImputa Age\nper titolo",          BLUE),
    ("CLEANER",       "KNN Imputer\nnum + cat\nFit su train",             TEAL),
    ("OUTLIERS",      "IQR + Z-score\nWinsorization\nSolo logreg",        PURPLE),
    ("MC HANDLER",    "Pearson |r|>0.85\nDrop feature\nSolo logreg",      RGBColor(0x6C, 0x35, 0x8A)),
    ("ENCODER",       "StandardScaler\nOne-Hot Enc.\nencode_target",      ORANGE),
    ("PCA",           "Varianza 95%\nSolo logreg\nRiduci dim.",           GREEN),
    ("MODEL",         "GridSearchCV\nbest_estimator_\nbest_score_",       RED),
]

bw = Inches(1.5)
bh = Inches(3.8)
sx = Inches(0.22)
gap = Inches(0.14)

for i, (name, desc, col) in enumerate(steps):
    bx = sx + i * (bw + gap)
    by = Inches(1.35)
    add_rect(sl, bx, by, bw, Inches(0.48), fill=col)
    txt(sl, name, bx, by + Inches(0.04), bw, Inches(0.42),
        size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, bx, by + Inches(0.48), bw, bh - Inches(0.48),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1))
    txt(sl, desc, bx + Inches(0.08), by + Inches(0.58),
        bw - Inches(0.14), bh - Inches(0.65), size=10.5, color=DARK_GRAY)
    if i < len(steps) - 1:
        ax = bx + bw + Inches(0.01)
        txt(sl, "->", ax, by + Inches(1.55), gap + Inches(0.1), Inches(0.5),
            size=16, bold=True, color=RGBColor(0xBB, 0xBB, 0xBB),
            align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.3), Inches(5.38), Inches(12.7), Inches(1.72),
         fill=RGBColor(0xFF, 0xFB, 0xF0), line=ORANGE, line_w=Pt(1))
txt(sl,
    "Regola anti-leakage: ogni transformer viene fittato ESCLUSIVAMENTE sui dati di train, "
    "poi applicato al test con transform().\n"
    "Encoder e PCAReducer hanno istanze separate per le due varianti del dataset "
    "(con/senza multicollinearita'), perche' le colonne presenti sono diverse "
    "e non si puo' riusare lo stesso scaler.",
    Inches(0.5), Inches(5.48), Inches(12.3), Inches(1.52),
    size=11.5, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# S7 — FEATURE ENGINEERING
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Feature Engineering -- Le Variabili Derivate",
       "7 nuove feature costruite sulla conoscenza del dominio Titanic")
pnum(sl, 7)

features = [
    ("FarePerPerson",
     ORANGE,
     "Fare / (SibSp + Parch + 1)",
     "Biglietti condivisi tra famiglie hanno tariffa unica. La tariffa pro-capite "
     "riflette la classe economica individuale, non il totale del gruppo.",
     "Fare grezzo rimosso per logreg/svm"),
    ("FamilySize",
     BLUE,
     "SibSp + Parch + 1",
     "Dimensione del nucleo familiare a bordo. Famiglie medie (2-4) avevano "
     "maggiore probabilita' di sopravvivenza rispetto a solitari o famiglie numerose.",
     "FamilySize=1 -> IsAlone=1"),
    ("IsAlone",
     TEAL,
     "(FamilySize == 1).astype(int)",
     "Flag per passeggeri senza familiari. I solitari avevano dinamiche di "
     "sopravvivenza diverse: meno vincoli ma meno supporto per le scialuppe.",
     "Feature booleana, basso costo"),
    ("CabinKnown",
     GREEN,
     "Cabin.notna().astype(int)",
     "Avere una cabina registrata era correlato con la classe sociale. Solo i "
     "passeggeri di 1a classe avevano generalmente la cabina assegnata.",
     "Converte l'assenza in informazione"),
    ("CabinDeck",
     PURPLE,
     "Cabin.str[0].fillna('U')",
     "Il ponte della cabina (A/B/C/D/E/F/G) determinava la distanza fisica "
     "dalle scialuppe. I ponti alti (A, B) erano piu' vicini al piano superiore.",
     "7 valori + U (unknown)"),
    ("SocialClass",
     RED,
     "Pclass x Sex  -> categoria",
     "Combina classe e genere: '1st_female', '3rd_male' ecc. Le donne di 1a "
     "classe sopravvivevano al 97%, gli uomini di 3a al 15%.",
     "Feature di interazione non lineare"),
    ("TitleGroup",
     NAVY,
     "regex r', ([^\\.]+)\\.'  su Name",
     "Il titolo (Mr, Mrs, Miss, Master, Rare) porta informazioni su sesso, eta' "
     "e status sociale. Master = bambino maschio, alta sopravvivenza.",
     "Usata anche per imputazione Age"),
]

for i, (name, col, formula, reason, note) in enumerate(features):
    cx = Inches(0.3)  if i % 2 == 0 else Inches(6.85)
    cy = Inches(1.3) + (i // 2) * Inches(1.55)
    if i == 6:
        cx = Inches(3.57)
    add_rect(sl, cx, cy, Inches(6.2), Inches(1.43),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1.5))
    add_rect(sl, cx, cy, Inches(6.2), Inches(0.37), fill=col)
    txt(sl, name, cx + Inches(0.12), cy + Inches(0.04),
        Inches(3.2), Inches(0.3), size=12, bold=True, color=WHITE)
    txt(sl, formula, cx + Inches(3.4), cy + Inches(0.06),
        Inches(2.65), Inches(0.27), size=9.5, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xCC), align=PP_ALIGN.RIGHT)
    txt(sl, reason, cx + Inches(0.12), cy + Inches(0.42),
        Inches(6.0), Inches(0.62), size=10.5, color=DARK_GRAY)
    add_rect(sl, cx + Inches(0.1), cy + Inches(1.07), Inches(6.0), Inches(0.3),
             fill=RGBColor(0xE8, 0xF4, 0xFF))
    txt(sl, f"  {note}", cx + Inches(0.18), cy + Inches(1.1),
        Inches(5.9), Inches(0.26), size=9.5, italic=True, color=NAVY)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32), fill=LIGHT_GRAY)
txt(sl,
    "PassengerId, Ticket e Cabin grezze vengono droppate: "
    "il loro contenuto utile e' gia' estratto nelle nuove variabili.",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S8 — IMPUTAZIONE AGE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Feature Engineering -- Imputazione Age per Titolo",
       "177 valori Age mancanti (19.9%) imputati con mediana per gruppo sociale")
pnum(sl, 8)

add_rect(sl, Inches(0.3), Inches(1.28), Inches(5.8), Inches(1.2),
         fill=RGBColor(0xFF, 0xF0, 0xF0), line=RED, line_w=Pt(1))
txt(sl, "Il problema: 177 Age mancanti su 891 (19.9%)",
    Inches(0.45), Inches(1.3), Inches(5.6), Inches(0.38),
    size=13, bold=True, color=RED)
txt(sl,
    "Imputare con la mediana globale (~28 anni) sbaglia sistematicamente: "
    "assegna 28 anni a 'Master Smith' (bambino di 3 anni) e a 'Mrs. Brown' (anziana).",
    Inches(0.45), Inches(1.7), Inches(5.6), Inches(0.68),
    size=11, color=DARK_GRAY)

add_rect(sl, Inches(0.3), Inches(2.62), Inches(5.8), Inches(4.56),
         fill=RGBColor(0xF0, 0xFF, 0xF0), line=GREEN, line_w=Pt(1.5))
add_rect(sl, Inches(0.3), Inches(2.62), Inches(5.8), Inches(0.4), fill=GREEN)
txt(sl, "La soluzione: mediana per TitleGroup (fit su train)",
    Inches(0.44), Inches(2.64), Inches(5.6), Inches(0.36),
    size=13, bold=True, color=WHITE)

groups = [
    ("Master", "~5 anni",  "Bambini maschi (figlio, nipote...)",   RGBColor(0x17, 0x8A, 0x8A)),
    ("Miss",   "~22 anni", "Donne non sposate, spesso giovani",    BLUE),
    ("Mrs",    "~36 anni", "Donne sposate / madri",                PURPLE),
    ("Mr",     "~30 anni", "Uomini adulti (la maggioranza)",       NAVY),
    ("Rare",   "mediana gruppo", "Rev, Dr, Col, Major, Lady...",   ORANGE),
]
for i, (grp, age, desc, col) in enumerate(groups):
    y = Inches(3.1) + i * Inches(0.77)
    add_rect(sl, Inches(0.45), y, Inches(1.0), Inches(0.52), fill=col)
    txt(sl, grp, Inches(0.45), y + Inches(0.1), Inches(1.0), Inches(0.35),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, age, Inches(1.55), y + Inches(0.1), Inches(1.5), Inches(0.35),
        size=13, bold=True, color=col)
    txt(sl, desc, Inches(3.15), y + Inches(0.1), Inches(2.75), Inches(0.35),
        size=10.5, color=DARK_GRAY)

add_rect(sl, Inches(6.4), Inches(1.28), Inches(6.5), Inches(5.9), fill=DARK_BG)
txt(sl, "feature_engineering.py", Inches(6.55), Inches(1.35),
    Inches(6.2), Inches(0.3), size=11, bold=True, color=ORANGE)
code = [
    ("_TITLE_MAP = {",                               TEAL),
    ("    'Master': 'Master',",                      WHITE),
    ("    'Miss': 'Miss', 'Ms': 'Miss',",            WHITE),
    ("    'Mrs': 'Mrs', 'Mme': 'Mrs',",              WHITE),
    ("    'Mr': 'Mr', 'Sir': 'Mr', ...",             WHITE),
    ("}  # tutti gli altri -> 'Rare'",               RGBColor(0x77, 0x99, 0x55)),
    ("",                                              WHITE),
    ("def fit(self, df):",                           TEAL),
    ("    titles = df['Name'].str.extract(",         WHITE),
    ("        r', ([^\\.]+)\\.')[0]",                WHITE),
    ("    grp = titles.map(_TITLE_MAP)",             WHITE),
    ("               .fillna('Rare')",               WHITE),
    ("    tmp = df.assign(TitleGroup=grp)",          WHITE),
    ("    self._age_medians = (",                    WHITE),
    ("        tmp.groupby('TitleGroup')",            WHITE),
    ("           ['Age'].median()",                  WHITE),
    ("           .to_dict()",                        WHITE),
    ("    )",                                         WHITE),
    ("    self._fallback = df['Age'].median()",      WHITE),
    ("",                                              WHITE),
    ("def transform(self, df):",                     TEAL),
    ("    # imputa PRIMA del cleaner",               RGBColor(0x77, 0x99, 0x55)),
    ("    mask = df['Age'].isna()",                  WHITE),
    ("    df.loc[mask, 'Age'] = grp[mask].map(",    WHITE),
    ("        self._age_medians",                    WHITE),
    ("    ).fillna(self._fallback)",                 WHITE),
]
for i, (line, col) in enumerate(code):
    txt(sl, line, Inches(6.55), Inches(1.72) + i * Inches(0.245),
        Inches(6.2), Inches(0.26), size=9.5, color=col)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32),
         fill=RGBColor(0xE8, 0xF4, 0xE8), line=GREEN, line_w=Pt(0.75))
txt(sl,
    "Fallback: se un gruppo non compare nel train (es. Rare vuoto), "
    "si usa la mediana globale del train -- nessun NaN residuo dopo la fase di fit.",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S9 — MODELLI
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Modelli & GridSearchCV",
       "5 algoritmi con iperparametri ottimizzati tramite cross-validation interna")
pnum(sl, 9)

models = [
    ("Logistic\nRegression",
     "Modello lineare, interpretabile.\nDataset: win+no-mc+no Fare+PCA",
     "C: [0.001, 0.01, 0.1, 1, 10]\npenalty: [l1, l2]\nsolver: [liblinear]",
     BLUE),
    ("Decision\nTree",
     "Albero decisionale.\nMassima interpretabilita'.",
     "max_depth: [3,5,7,10]\nmin_samples_split: [2,5,10]\ncriterion: [gini, entropy]",
     TEAL),
    ("Random\nForest",
     "Ensemble di alberi.\nn_jobs=1 sull'estimator\n(fix deadlock Docker).",
     "n_estimators: [100,200]\nmax_depth: [5,10]\nmin_samples_split: [2,5]",
     GREEN),
    ("XGBoost",
     "Gradient boosting sequenziale.\nnthread=1 sull'estimator\n(fix deadlock Docker).",
     "n_estimators: [100,200]\nmax_depth: [3,5]\nlearning_rate: [0.05,0.1]",
     ORANGE),
    ("SVM",
     "Support Vector Machine.\nprobability=True per\nthreshold tuning.",
     "kernel: [rbf, linear, poly]\nC: [0.1, 1, 10, 100]\ngamma: [scale, auto]",
     RED),
]

for i, (name, desc, params, col) in enumerate(models):
    bx = Inches(0.28) + i * Inches(2.6)
    add_rect(sl, bx, Inches(1.28), Inches(2.45), Inches(0.55), fill=col)
    txt(sl, name, bx, Inches(1.3), Inches(2.45), Inches(0.52),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, bx, Inches(1.83), Inches(2.45), Inches(2.6),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1))
    txt(sl, desc, bx + Inches(0.1), Inches(1.9),
        Inches(2.28), Inches(1.4), size=10.5, color=DARK_GRAY)
    add_rect(sl, bx, Inches(4.43), Inches(2.45), Inches(1.7),
             fill=RGBColor(0xE8, 0xF4, 0xFF), line=col, line_w=Pt(0.75))
    txt(sl, "Param grid:", bx + Inches(0.1), Inches(4.5),
        Inches(2.28), Inches(0.3), size=10.5, bold=True, color=col)
    txt(sl, params, bx + Inches(0.1), Inches(4.82),
        Inches(2.28), Inches(1.2), size=10, color=DARK_GRAY)

add_rect(sl, Inches(0.28), Inches(6.25), Inches(12.77), Inches(1.12),
         fill=RGBColor(0xE8, 0xF4, 0xE8), line=GREEN, line_w=Pt(1))
txt(sl,
    "GridSearchCV: cv=5, scoring='f1_weighted', n_jobs=-1 (parallelo a livello di fold).\n"
    "best_score_ e' la stima onesta della cross-validation interna, non calcolata sul train.\n"
    "best_estimator_ viene poi ri-valutato con cross_val_score(scoring='accuracy') per ottenere la stima Kaggle.",
    Inches(0.5), Inches(6.32), Inches(12.4), Inches(1.0),
    size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# S10 — TECNICHE AVANZATE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Tecniche Avanzate",
       "PCA  |  Multicollinearity  |  Threshold Tuning  |  Model Caching")
pnum(sl, 10)

techs = [
    ("PCA -- Principal Component Analysis", BLUE,
     ["Mantiene il 95% della varianza spiegata",
      "Solo per logreg: beneficia di feature ortogonali",
      "PCAReducer separato per ogni modello nell'apply_to",
      "n_components=0.95: numero componenti adattivo"],
     "LogReg e' sensibile a feature correlate: "
     "PCA le ortogonalizza riducendo instabilita' numerica e migliorando la convergenza."),
    ("Multicollinearity Handler", TEAL,
     ["Correlazione di Pearson tra coppie di feature numeriche",
      "Rimuove automaticamente se |r| > soglia (0.85)",
      "encoder_mc: istanza separata per il dataset ridotto",
      "Applicato prima dell'encoding (opera su valori grezzi)"],
     "Feature ridondanti non aggiungono informazione ma aumentano la varianza "
     "degli stimatori lineari e complicano l'interpretazione dei coefficienti."),
    ("Decision Threshold Tuning", ORANGE,
     ["Default soglia: 0.5 (P > 0.5 -> predicted=1)",
      "Scannerizza soglie da 0.30 a 0.70 con passo 0.01",
      "Usa cross_val_predict(method='predict_proba') -> OOF probs",
      "Ottimizza F1 (config: accuracy, recall, precision)"],
     "Con classi sbilanciate (38% sopravvissuti) la soglia 0.5 non e' ottimale "
     "per F1. Spostarla verso 0.40 recupera recall sui positivi."),
    ("Model Caching", GREEN,
     ["Salva best_estimator_ su .pkl dopo il training",
      "Hash MD5 di loader+models+cleaner+outliers+encoder+pca",
      "Config cambiata -> hash diverso -> retraining automatico",
      "Disabilitato durante sviluppo (enabled: false)"],
     "GridSearchCV su 5 modelli richiede diversi minuti. Con caching, "
     "run successive con la stessa configurazione caricano il modello in millisecondi."),
]

for i, (title, col, points, why) in enumerate(techs):
    row, c = i // 2, i % 2
    bx = Inches(0.3) + c * Inches(6.55)
    by = Inches(1.28) + row * Inches(3.0)
    add_rect(sl, bx, by, Inches(6.25), Inches(2.82),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1.5))
    add_rect(sl, bx, by, Inches(6.25), Inches(0.4), fill=col)
    txt(sl, title, bx + Inches(0.12), by + Inches(0.04),
        Inches(6.05), Inches(0.34), size=13, bold=True, color=WHITE)
    for j, pt in enumerate(points):
        txt(sl, f"  {pt}", bx + Inches(0.14), by + Inches(0.5) + j * Inches(0.33),
            Inches(6.0), Inches(0.3), size=11, color=DARK_GRAY)
    add_rect(sl, bx + Inches(0.1), by + Inches(1.9), Inches(6.05), Inches(0.8),
             fill=RGBColor(0xFF, 0xFB, 0xF0), line=col, line_w=Pt(0.5))
    txt(sl, f"Perche': {why}", bx + Inches(0.2), by + Inches(1.96),
        Inches(5.88), Inches(0.72), size=10, italic=True, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# S11 — PARALLELISMO & DOCKER
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Parallelismo & Docker -- Fix Deadlock",
       "Il problema del nested parallelism e come risolverlo in container")
pnum(sl, 11)

add_rect(sl, Inches(0.3), Inches(1.28), Inches(12.7), Inches(1.65),
         fill=RGBColor(0xFF, 0xF0, 0xF0), line=RED, line_w=Pt(1.5))
add_rect(sl, Inches(0.3), Inches(1.28), Inches(12.7), Inches(0.4), fill=RED)
txt(sl, "Il Problema: nested joblib parallelism in Docker",
    Inches(0.45), Inches(1.3), Inches(12.4), Inches(0.36),
    size=14, bold=True, color=WHITE)
txt(sl,
    "GridSearchCV(n_jobs=-1) usa il backend loky (processi) per parallelizzare i fold. "
    "Se l'estimator usa anche n_jobs=-1 (RF) o nthread>1 (XGBoost), "
    "Docker cerca di forkare sotto-processi dentro sotto-processi: deadlock o OOM. "
    "Il processo si blocca silenziosamente senza messaggi di errore.",
    Inches(0.45), Inches(1.72), Inches(12.4), Inches(1.1),
    size=11.5, color=DARK_GRAY)

fixes = [
    ("GridSearchCV",             BLUE,   "n_jobs=-1",
     "Parallelizza i fold di cross-validation.\nSicuro: ogni fold e' un processo "
     "indipendente a livello di GridSearchCV."),
    ("RandomForestClassifier",   GREEN,  "n_jobs=1",
     "RF e' gia' un ensemble interno di alberi.\nSenza n_jobs=-1 non forka processi "
     "dentro i processi dei fold."),
    ("XGBClassifier",            ORANGE, "nthread=1",
     "XGBoost usa il suo thread pool interno.\nnthread=1 forza single-thread "
     "evitando conflitti con il loky pool."),
    ("ENV LOKY_MAX_CPU_COUNT",   TEAL,   "=8",
     "Variabile nel Dockerfile.\nLoky tenta di rilevare i core del host, "
     "che in Docker e' instabile. Impostarla stabilizza i worker."),
]

for i, (component, col, value, desc) in enumerate(fixes):
    bx = Inches(0.3) + i * Inches(3.2)
    by = Inches(3.1)
    add_rect(sl, bx, by, Inches(3.0), Inches(3.9),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1.5))
    add_rect(sl, bx, by, Inches(3.0), Inches(0.4), fill=col)
    txt(sl, component, bx + Inches(0.1), by + Inches(0.04),
        Inches(2.82), Inches(0.34), size=12, bold=True, color=WHITE)
    add_rect(sl, bx + Inches(0.5), by + Inches(0.5), Inches(2.0), Inches(0.52),
             fill=DARK_BG)
    txt(sl, value, bx + Inches(0.5), by + Inches(0.56),
        Inches(2.0), Inches(0.38), size=16, bold=True, color=ORANGE,
        align=PP_ALIGN.CENTER)
    txt(sl, desc, bx + Inches(0.12), by + Inches(1.15),
        Inches(2.78), Inches(2.65), size=10.5, color=DARK_GRAY)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32),
         fill=RGBColor(0xE8, 0xF4, 0xE8), line=GREEN, line_w=Pt(0.75))
txt(sl,
    "ENV PYTHONUNBUFFERED=1 nel Dockerfile: senza di essa Python bufferizza stdout "
    "e non stampa nulla durante il training nel container.",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S12 — VALUTAZIONE LOCALE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Valutazione Locale -- evaluate_local.py",
       "Come stimare l'accuracy Kaggle senza fare upload")
pnum(sl, 12)

add_rect(sl, Inches(0.3), Inches(1.28), Inches(12.7), Inches(0.9),
         fill=RGBColor(0xFF, 0xF8, 0xE8), line=ORANGE, line_w=Pt(1))
txt(sl,
    "I label del test set (passeggeri 892-1309) sono segreti -- Kaggle non li pubblica. "
    "La valutazione onesta si fa sul train set (891 campioni) con cross-validation. "
    "evaluate_local.py implementa lo stesso approccio del notebook di riferimento.",
    Inches(0.5), Inches(1.36), Inches(12.3), Inches(0.74),
    size=12, color=DARK_GRAY)

add_rect(sl, Inches(0.3), Inches(2.3), Inches(5.9), Inches(3.6),
         fill=RGBColor(0xFF, 0xF0, 0xF0), line=RED, line_w=Pt(1.5))
add_rect(sl, Inches(0.3), Inches(2.3), Inches(5.9), Inches(0.4), fill=RED)
txt(sl, "Approccio 1 -- Holdout 80/20",
    Inches(0.44), Inches(2.32), Inches(5.7), Inches(0.36),
    size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "StratifiedShuffleSplit(test_size=0.2, seed=42)",
    "Preprocessing fittato SOLO sul train split (713)",
    "Valutazione sul val split (178 righe)",
    "Pro: nessun leakage del preprocessing",
    "Contro: stima variabile (dipende dal seed)",
    "Utile per: debug e verifiche singole",
]):
    txt(sl, f"  {item}", Inches(0.44), Inches(2.82) + i * Inches(0.4),
        Inches(5.7), Inches(0.36), size=11, color=DARK_GRAY)

add_rect(sl, Inches(6.6), Inches(2.3), Inches(6.1), Inches(3.6),
         fill=RGBColor(0xF0, 0xFF, 0xF0), line=GREEN, line_w=Pt(1.5))
add_rect(sl, Inches(6.6), Inches(2.3), Inches(6.1), Inches(0.4), fill=GREEN)
txt(sl, "Approccio 2 -- StratifiedKFold(5)  * Usato",
    Inches(6.74), Inches(2.32), Inches(5.9), Inches(0.36),
    size=13, bold=True, color=WHITE)
for i, item in enumerate([
    "Preprocessing su tutti i 891 campioni",
    "GridSearchCV -> best_estimator_",
    "cross_val_score(best_estimator, X_enc, y, cv=5)",
    "Pro: usa tutti i dati, stima piu' stabile",
    "Output: mean accuracy +/- std tra i fold",
    "Stesso approccio del notebook di riferimento",
]):
    txt(sl, f"  {item}", Inches(6.74), Inches(2.82) + i * Inches(0.4),
        Inches(5.9), Inches(0.36), size=11, color=DARK_GRAY)

add_rect(sl, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.15), fill=DARK_BG)
txt(sl, "evaluate_local.py", Inches(0.5), Inches(6.1),
    Inches(4.0), Inches(0.3), size=11, bold=True, color=ORANGE)
for i, (line, col) in enumerate([
    ("cv = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)",                         TEAL),
    ("scores = cross_val_score(model.model, X_enc, y_enc, cv=cv, scoring='accuracy', n_jobs=-1)", WHITE),
    ("print(f'Mean: {scores.mean():.4f}  Std: {scores.std():.4f}')",                            RGBColor(0xAA, 0xCC, 0xFF)),
]):
    txt(sl, line, Inches(0.5), Inches(6.45) + i * Inches(0.26),
        Inches(12.5), Inches(0.26), size=10, color=col)


# ═══════════════════════════════════════════════════════════════════
# S13 — RISULTATI
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Risultati -- Performance CV Locale",
       "StratifiedKFold(5), scoring='accuracy' su train.csv (891 campioni)")
pnum(sl, 13)

results = [
    ("Random Forest",        "orig",                    0.8440, 0.0128, GREEN),
    ("XGBoost",              "orig",                    0.8361, 0.0064, ORANGE),
    ("SVM",                  "orig (no Fare)",           0.8328, 0.0156, RED),
    ("Logistic Regression",  "win+no-mc (no Fare)",      0.8305, 0.0072, BLUE),
    ("Decision Tree",        "orig",                    0.8260, 0.0164, TEAL),
    ("LogReg + PCA",         "win+no-mc (no Fare)+PCA",  0.8148, 0.0096, PURPLE),
    ("Baseline (tutti 0)",   "--",                       0.6162, 0.0,    RGBColor(0xAA, 0xAA, 0xAA)),
]

add_rect(sl, Inches(0.3), Inches(1.28), Inches(7.5), Inches(0.42), fill=NAVY)
for j, (hdr, cx) in enumerate(zip(
    ["Modello", "Dataset", "Accuracy", "+/- Std", "Bar"],
    [Inches(0.32), Inches(2.16), Inches(4.4), Inches(5.38), Inches(6.27)]
)):
    txt(sl, hdr, cx + Inches(0.06), Inches(1.3), Inches(1.8), Inches(0.35),
        size=11, bold=True, color=WHITE)

for i, (name, dataset, acc, std, col) in enumerate(results):
    y = Inches(1.7) + i * Inches(0.62)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(sl, Inches(0.3), y, Inches(7.5), Inches(0.6), fill=bg,
             line=RGBColor(0xDD, 0xDD, 0xDD), line_w=Pt(0.5))
    txt(sl, name,    Inches(0.38), y + Inches(0.14), Inches(1.8), Inches(0.32),
        size=11, bold=(i == 0), color=(col if i == 0 else DARK_GRAY))
    txt(sl, dataset, Inches(2.22), y + Inches(0.14), Inches(2.15), Inches(0.32),
        size=10, color=DARK_GRAY)
    txt(sl, f"{acc:.4f}", Inches(4.46), y + Inches(0.14), Inches(0.9), Inches(0.32),
        size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)
    txt(sl, f"+-{std:.4f}", Inches(5.44), y + Inches(0.14), Inches(0.8), Inches(0.32),
        size=10, color=DARK_GRAY)
    if acc > 0.65:
        bar_w = (acc - 0.61) / 0.26 * Inches(1.5)
        add_rect(sl, Inches(6.33), y + Inches(0.2), bar_w, Inches(0.22), fill=col)

add_rect(sl, Inches(0.3), Inches(6.05), Inches(7.5), Inches(0.5),
         fill=RGBColor(0xE8, 0xF4, 0xE8), line=GREEN, line_w=Pt(1))
txt(sl, "  Best: Random Forest 84.4%  |  +22.8 pp sulla baseline naive  |  "
    "Spread modelli: 2.9 pp tra il migliore e il peggiore (esclusa baseline)",
    Inches(0.42), Inches(6.12), Inches(7.2), Inches(0.35),
    size=11, bold=True, color=DARK_GRAY)

add_rect(sl, Inches(8.1), Inches(1.28), Inches(4.9), Inches(5.28),
         fill=RGBColor(0xF8, 0xF9, 0xFF), line=NAVY, line_w=Pt(1.5))
add_rect(sl, Inches(8.1), Inches(1.28), Inches(4.9), Inches(0.42), fill=NAVY)
txt(sl, "Analisi dei risultati",
    Inches(8.22), Inches(1.3), Inches(4.7), Inches(0.36),
    size=13, bold=True, color=WHITE)
observations = [
    ("Random Forest 84.4%",
     "Robusto all'overfitting. Gestisce bene outlier e feature categoriche del Titanic."),
    ("XGBoost vs RF (0.8%)",
     "Differenza contenuta: con ensemble i due si complementano."),
    ("SVM competitivo 83.3%",
     "Kernel RBF cattura bene confini non lineari del dataset."),
    ("LogReg + PCA peggiore",
     "PCA perde informazione su questo dataset piccolo (891 righe)."),
    ("Std RF 0.0128",
     "Deviazione standard accettabile: il modello e' stabile tra i fold."),
    ("Gap vs notebook",
     "Il notebook raggiunge 84% con VotingClassifier ensemble."),
]
for i, (title, desc) in enumerate(observations):
    y = Inches(1.82) + i * Inches(0.75)
    txt(sl, title, Inches(8.22), y, Inches(4.6), Inches(0.28),
        size=11, bold=True, color=BLUE)
    txt(sl, desc, Inches(8.22), y + Inches(0.3), Inches(4.6), Inches(0.38),
        size=10, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# S14 — CONFRONTO FUNGHI
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Confronto: Progetto Funghi -> Titanic",
       "Evoluzione dell'architettura rispetto al progetto precedente")
pnum(sl, 14)

headers_t = ["Aspetto", "Progetto Funghi", "Progetto Titanic", "Impatto"]
col_ws = [Inches(2.1), Inches(3.3), Inches(4.7), Inches(2.6)]
col_xs = [Inches(0.3), Inches(2.44), Inches(5.78), Inches(10.52)]

for j, (hdr, cw, cx) in enumerate(zip(headers_t, col_ws, col_xs)):
    add_rect(sl, cx, Inches(1.28), cw, Inches(0.42), fill=NAVY)
    txt(sl, hdr, cx + Inches(0.08), Inches(1.3), cw - Inches(0.12), Inches(0.35),
        size=12, bold=True, color=WHITE)

rows_data = [
    ("Containerizzazione",  "Script locale Python",
     "Docker (python:3.11-slim)\nPYTHONPATH, LOKY, UNBUFFERED",
     "Riproducibile ovunque", GREEN),
    ("API REST",            "Nessuna",
     "FastAPI: /train /evaluate /predict\nRisposta JSON + submission CSV",
     "Integrabile in app", GREEN),
    ("Libreria ML",         "Classi non strutturate",
     "AbstractModel, GenericCleaner...\nconfigure -> fit -> transform",
     "Riusabile su altri dataset", GREEN),
    ("Feature Engineering", "Minimal o nessuna",
     "7 nuove feature derivate\nImputazione Age per TitleGroup",
     "Piu' info al modello", GREEN),
    ("Preprocessing",       "Cleaner base",
     "Outlier + Winsorization\nMulticollinearity handler\nEncoder separati per variante",
     "Pipeline coerente", GREEN),
    ("Riduzione dim.",      "Nessuna",
     "PCA 95% varianza\nPCAReducer per ogni modello\nEvita feature mismatch",
     "LogReg piu' stabile", TEAL),
    ("Parallelismo",        "Default joblib",
     "Fix deadlock Docker\nn_jobs=1 estimator, n_jobs=-1 GridSearchCV",
     "Nessun hang", ORANGE),
    ("Valutazione",         "Train accuracy",
     "StratifiedKFold(5) + cross_val_score\nLocal eval script dedicato",
     "Stima Kaggle onesta", GREEN),
]

for i, (aspect, old, new, impact, imp_col) in enumerate(rows_data):
    y = Inches(1.7) + i * Inches(0.65)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (text, cw, cx) in enumerate(zip([aspect, old, new, impact], col_ws, col_xs)):
        add_rect(sl, cx, y, cw, Inches(0.63), fill=bg,
                 line=RGBColor(0xDD, 0xDD, 0xDD), line_w=Pt(0.5))
        col_text = (RED if j == 1 else (imp_col if j == 3 else DARK_GRAY))
        txt(sl, text, cx + Inches(0.08), y + Inches(0.06),
            cw - Inches(0.12), Inches(0.54), size=10,
            color=col_text, bold=(j == 3))

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32), fill=LIGHT_GRAY)
txt(sl,
    "Da script monolitico a sistema modulare, containerizzato, con API, "
    "valutazione onesta e libreria ML riutilizzabile su qualsiasi dataset tabellare.",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S15 — PROSSIMI PASSI
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Prossimi Passi -- Ensemble & Roadmap",
       "Come superare l'84% e raggiungere le performance dei notebook di riferimento")
pnum(sl, 15)

next_steps = [
    ("Soft Voting Ensemble", ORANGE,
     "VotingClassifier(voting='soft') su logreg + xgboost + random_forest.\n"
     "Media delle probabilita' dei 3 modelli migliori.\n"
     "Il notebook di riferimento raggiunge ~84% con questo approccio.\n"
     "Implementazione: 5 righe con sklearn VotingClassifier.",
     "Guadagno atteso: +1-2 pp su 84.4%"),
    ("CV-Bagging su Random Forest", GREEN,
     "Approccio dell'articolo leaderboard (~87%).\n"
     "Training del modello su ogni fold di StratifiedKFold(5).\n"
     "Predizioni test = media delle probabilita' dei 5 modelli.\n"
     "RF 1750 alberi, max_depth=7 -> piu' espressivo del nostro.",
     "Guadagno atteso: +2-3 pp, submission piu' stabile"),
    ("Stacking con Meta-learner", BLUE,
     "Modelli base: RF, XGBoost, SVM -> predizioni OOF come feature.\n"
     "Meta-learner: LogisticRegression sulle predizioni base.\n"
     "Attenzione al leakage: base e meta su split separati.\n"
     "Implementazione con sklearn StackingClassifier.",
     "Guadagno atteso: +1-2 pp, complessita' maggiore"),
    ("Feature Interaction", TEAL,
     "Sex x Pclass: feature di interazione esplicita.\n"
     "TicketFreq: quante persone condividono il ticket.\n"
     "AgeGroup: fasce eta' (Child <12, Teen, Adult, Senior).\n"
     "Fare binning: fasce tariffarie invece di valore continuo.",
     "Guadagno atteso: +0.5-1 pp"),
]

for i, (title, col, desc, expected) in enumerate(next_steps):
    row, c = i // 2, i % 2
    bx = Inches(0.3) + c * Inches(6.55)
    by = Inches(1.28) + row * Inches(2.98)
    add_rect(sl, bx, by, Inches(6.2), Inches(2.8),
             fill=LIGHT_GRAY, line=col, line_w=Pt(1.5))
    add_rect(sl, bx, by, Inches(6.2), Inches(0.4), fill=col)
    txt(sl, title, bx + Inches(0.12), by + Inches(0.04),
        Inches(6.0), Inches(0.34), size=13, bold=True, color=WHITE)
    txt(sl, desc, bx + Inches(0.12), by + Inches(0.48),
        Inches(5.97), Inches(1.52), size=10.5, color=DARK_GRAY)
    add_rect(sl, bx + Inches(0.1), by + Inches(2.08), Inches(6.0), Inches(0.6),
             fill=RGBColor(0xFF, 0xFB, 0xF0), line=col, line_w=Pt(0.5))
    txt(sl, f"  {expected}", bx + Inches(0.2), by + Inches(2.14),
        Inches(5.85), Inches(0.5), size=10.5, bold=True, color=col)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32), fill=LIGHT_GRAY)
txt(sl,
    "Benchmark: Notebook VotingClassifier ~84%  |  "
    "Articolo CV-Bagging RF ~87%  |  Top Kaggle >90% (deep feature engineering + ensembles)",
    Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.26),
    size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# S16 — FINE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(3.2), W, Inches(0.08), fill=ORANGE)
add_rect(sl, 0, Inches(4.55), W, Inches(0.08), fill=TEAL)

txt(sl, "Grazie", Inches(0.5), Inches(1.0), W - Inches(1.0), Inches(1.2),
    size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Titanic Survival Prediction -- ML Pipeline Avanzata con Libreria Generalizzata",
    Inches(0.5), Inches(2.15), W - Inches(1.0), Inches(0.6),
    size=19, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER, italic=True)

tech_badges = [
    ("Docker",            "Containerizzazione"),
    ("FastAPI",           "REST API"),
    ("AbstractModel",     "Libreria generalizzata"),
    ("GridSearchCV",      "Ottimizzazione HP"),
    ("PCA + MC Handler",  "Riduzione dim."),
    ("StratifiedKFold",   "Valutazione onesta"),
    ("XGBoost 84.4%",     "Best model"),
    ("Feature Eng. x7",   "Nuove variabili"),
]
for i, (tech, desc) in enumerate(tech_badges):
    col_n = i % 4
    row_n = i // 4
    bx = Inches(0.9) + col_n * Inches(2.9)
    by = Inches(4.85) + row_n * Inches(1.08)
    add_rect(sl, bx, by, Inches(2.6), Inches(0.88),
             fill=RGBColor(0x2A, 0x3A, 0x6E), line=TEAL, line_w=Pt(1))
    txt(sl, tech, bx, by + Inches(0.06), Inches(2.6), Inches(0.4),
        size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx, by + Inches(0.46), Inches(2.6), Inches(0.35),
        size=10, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

txt(sl, "Corso Data Analytics  |  2025/26  |  Andrea Morlotti",
    Inches(0.5), H - Inches(0.48), W - Inches(1.0), Inches(0.35),
    size=11, color=RGBColor(0x55, 0x66, 0x88), align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save('/output/titanic_presentation.pptx')
print(f"Salvata: /output/titanic_presentation.pptx  ({len(prs.slides)} slide)")
